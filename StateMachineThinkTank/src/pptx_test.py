from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def html_to_slides(html_content, output_file='presentation.pptx'):
    """
    Convert HTML content to PowerPoint slides.
    
    Args:
        html_content (str): HTML content as string
        output_file (str): Output PowerPoint file name
    """
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Parse HTML content
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Function to add a slide with title and content
    def add_content_slide(title, content):
        slide_layout = prs.slide_layouts[1]  # Using layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Set content
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.text = content
        text_frame.paragraphs[0].font.size = Pt(24)
        
        return slide

    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = soup.find('h1').text
    
    # Process each section (h2) as a separate slide
    sections = soup.find_all('h2')
    for section in sections:
        # Get the title
        slide_title = section.text
        
        # Get the content following this h2 until the next h2
        content = []
        current = section.next_sibling
        
        while current and current.name != 'h2':
            if current.name == 'h3':
                content.append(f"\n{current.text}")
            elif current.name == 'p':
                content.append(current.text)
            current = current.next_sibling
        
        # Join all content with proper spacing
        slide_content = '\n\n'.join(content)
        
        # Create the slide
        add_content_slide(slide_title, slide_content)
    
    # Save the presentation
    prs.save(output_file)

def main():
    # Read HTML file
    html_text = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Executive Summary - State-Driven AI Agent</title>
    <style>
        body {
            font-family: Arial, sans-serif;
            color: #333;
            line-height: 1.6;
        }
        h1, h2 {
            color: #2E3A59;
        }
        h1 {
            font-size: 2em;
            text-transform: uppercase;
            border-bottom: 2px solid #E2E8F0;
            padding-bottom: 10px;
        }
        h2 {
            font-size: 1.5em;
            margin-top: 1.5em;
            color: #4A5568;
        }
        .section {
            margin-top: 2em;
            border-top: 1px solid #CBD5E0;
            padding-top: 1em;
        }
        p {
            margin-top: 0.5em;
            color: #4A5568;
        }
    </style>
</head>
<body>

    <!-- Main Title Slide -->
    <h1>State-Driven AI Agent Framework</h1>
    <p>A modular and adaptive model for interactive, scalable AI applications</p>

    <!-- Introduction Slide -->
    <div class="section">
        <h2>Introduction</h2>
        <p>This state-driven AI framework is designed for high adaptability, enabling personalized, scalable interactions across diverse applications. Through states like Initialization, Session Types, Interactive Loop, and Finalization, the AI can dynamically guide users in complex environments, offering a tailored and efficient experience.</p>
    </div>

    <!-- Product Overview Slide (Technical Architecture) -->
    <div class="section">
        <h2>Technical Architecture</h2>
        <h3>Michael Taylor, Technical Architect</h3>
        <p>The framework’s state-based architecture provides a modular, extensible structure that supports scalability and flexibility. Its clear, defined pathways ensure the AI agent can adapt seamlessly to varied workflows, making it suitable for applications in customer support and project management where reliability and adaptability are crucial.</p>
    </div>

    <!-- Market Comparison Slide (Project Management) -->
    <div class="section">
        <h2>Project Management Benefits</h2>
        <h3>Kevin Fears, Project Manager</h3>
        <p>This structured AI model enhances user experiences, particularly in onboarding and guided workflows. It’s ideal for enterprise environments, enabling the AI to guide users in alignment with project objectives. Key applications include IT deployments and support scenarios, where the AI can manage common troubleshooting and user training tasks efficiently.</p>
    </div>

    <!-- Market Overview Slide (Business Value) -->
    <div class="section">
        <h2>Business Value and Scalability</h2>
        <h3>Sarah Williams, Venture Capitalist (Non-Technical)</h3>
        <p>This model enables affordable scalability by automating repetitive customer interactions, allowing companies to provide quality support without increasing headcount. In practice, it’s a strong fit for SaaS platforms and advisory services, offering guidance and customization at scale, ideal for rapid-growth environments.</p>
    </div>

    <!-- Growth Strategy Slide (Emerging Markets) -->
    <div class="section">
        <h2>Market Expansion in Emerging Sectors</h2>
        <h3>Rachel Bennett, Venture Capitalist (Emerging Markets)</h3>
        <p>The state-driven approach supports companies entering traditional or emerging markets, where it can guide users through complex processes and mitigate adoption barriers. The model is ideal for sectors like finance or healthcare, which require user education and gradual market entry strategies.</p>
    </div>

    <!-- Ethics and User Interaction Quality -->
    <div class="section">
        <h2>Ethics and User Interaction Quality</h2>
        <h3>Dr. Elena Chavez, AI Ethics Specialist</h3>
        <p>The Interactive Loop allows the AI to adjust in real-time based on user feedback, helping manage biases and maintaining ethical transparency. This is especially important for applications in education or healthcare, where unbiased guidance is crucial. This adaptive design ensures that interactions remain aligned with user expectations and psychological well-being.</p>
    </div>

    <!-- Sales Potential and Customer Engagement -->
    <div class="section">
        <h2>Sales Potential and Customer Engagement</h2>
        <h3>David Brooks, Sales Specialist</h3>
        <p>The model’s adaptability is highly suited to sales, where the AI can guide potential clients through relevant information based on their responses. This real-time customization enhances customer engagement and can significantly improve lead qualification and conversion rates in B2B environments.</p>
    </div>

    <!-- Technical Feasibility and Innovation -->
    <div class="section">
        <h2>Technical Feasibility and Innovation</h2>
        <h3>John Smith, Technical Venture Capitalist</h3>
        <p>The modular architecture allows for scalability and ease of innovation, making it a strong fit for high-growth startups. This model supports technical consulting, decision-support tools, and any environment where user interactions evolve continuously, making it an ideal choice for investor-backed growth.</p>
    </div>

    <!-- Customer Experience and Retention -->
    <div class="section">
        <h2>Customer Experience and Retention</h2>
        <h3>Jessica Lee, Account Executive</h3>
        <p>The User Feedback Prompt enables a personalized experience by allowing users to control interaction flow. This reduces friction, supporting customer retention and satisfaction, especially in SaaS models where customer success drives long-term engagement and account expansion opportunities.</p>
    </div>

</body>
</html>

"""
    # Convert to slides
    html_to_slides(html_text, 'executive_summary_presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    main()